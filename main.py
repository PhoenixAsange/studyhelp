from ctypes import wstring_at

from pptx import Presentation
from langchain_ollama import OllamaLLM
from langchain_core.prompts import ChatPromptTemplate
import os

template = """"
    You are a highly intelligent and helpful study assistant. You help students by answering questions based on the provided course material.

    User course content (reference when answering questions): {content}

    Conversation history:
    {context}

    Question:
    {question}
    
    Answer:
    """
model = OllamaLLM(model="llama3")
prompt = ChatPromptTemplate.from_template(template)
chain = prompt | model

def extractText(filePath):
    file = Presentation(filePath)
    text = []
    for slide in file.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    text = "\n".join(text)
    print(f"\nExtracted Text (Preview):\n {text[:50]}...")
    return text

def modelConversation():
    context = ""
    print("Type 'exit' to exit conversation.")
    while True:
        print("Bot: What would you like help with?")
        userInput = input("You: ")
        if userInput.lower() == "exit":
            break

        result = chain.invoke({"content": content, "context": context, "question": userInput})
        print("Bot:", result)
        context += f"\nUser: {userInput}\nAI: {result}"

while True:
    filePath = input("Enter file path (single = file, multiple = directory): ")
    try:
        if os.path.isdir(filePath):
            content = []
            fileAmount = 0
            for file in os.listdir(filePath):
                fullPath = os.path.join(filePath, file)
                content.append(extractText(fullPath))
                fileAmount += 1
            print(f"{fileAmount} files have been parsed")
        else:
            content = extractText(filePath)
        break
    except Exception as e:
        print(f"Error: {e}. Please ensure the directory/file path is correct.")

if __name__ == "__main__":
    modelConversation()